import os
import re
import json
import asyncio
from typing import List, Dict, Tuple
import argparse
from dotenv import load_dotenv
from googleapiclient.discovery import build
from tenacity import (
    retry,
    stop_after_attempt,
    wait_exponential_jitter,
    retry_if_exception_type,
)
import pptx
import PyPDF2

# --- third‑party async client (openai>=1.14.0) ---
from openai import AsyncOpenAI, OpenAIError
import sys


load_dotenv()

# ---------------------------------------------------------------------------
# ENV / GLOBALS
# ---------------------------------------------------------------------------
client = AsyncOpenAI()
GOOGLE_API_KEY = os.getenv("GOOGLE_API_KEY")
CUSTOM_SEARCH_ENGINE_ID = os.getenv("CUSTOM_SEARCH_ENGINE_ID")
MAX_SOURCES = 5  # 上位何件を証拠にするか

# ---------------------------------------------------------------------------
# UTILS
# ---------------------------------------------------------------------------

def read_document(file_path):
    print(f"'{file_path}' を読み込んでいます...")
    _, extension = os.path.splitext(file_path)
    text = ""
    try:
        if extension == ".pptx":
            prs = pptx.Presentation(file_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
        elif extension == ".pdf":
            with open(file_path, "rb") as f:
                reader = PyPDF2.PdfReader(f)
                for page in reader.pages:
                    text += page.extract_text() + "\n"
        else:
            return None, f"対応していないファイル形式です: {extension}"
        print("ファイルの読み込みが完了しました。")
        return text, None
    except Exception as e:
        return None, f"ファイルの読み込み中にエラーが発生しました: {e}"
    
def extract_lines(slide):
    """(text, level, y_pos, x_pos) を返す"""
    lines = []
    for shape in slide.shapes:
        if not hasattr(shape, "text_frame"):  # 図形・画像などは無視
            continue
        for p in shape.text_frame.paragraphs:
            txt = p.text.strip()
            if not txt:
                continue
            # 箇条書きレベル: 0=タイトル/本文、1=•、2=–– など
            lvl = p.level
            y, x, _, _ = shape.left, shape.top, shape.width, shape.height
            lines.append((txt, lvl, y, x))
    # 画面上: 上→下、左→右
    return sorted(lines, key=lambda t: (t[2], t[3]))


def sanitize_text(text: str) -> str:
    """余分な空白/改行を削除"""
    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r"\n{2,}", "\n", text)
    return text.strip()
    


def get_source_reliability(url: str) -> Tuple[str, int]:
    """ドメインごとに信頼度ラベルとスコアを返す"""
    url = url.lower()
    high = [
        ".sciencemag.org",
        "nature.com",
        "cell.com",
        "arxiv.org",
        "ieee.org",
        "acm.org",
        "pubmed.ncbi.nlm.nih.gov",
        ".go.jp",
        ".gov",
        ".ac.jp",
        ".edu",
    ]
    med = [
        "wikipedia.org",
        "bbc.com",
        "reuters.com",
        "apnews.com",
        "nytimes.com",
        "nikkei.com",
        "asahi.com",
        "yomiuri.co.jp",
        "stat.go.jp",
    ]
    if any(d in url for d in high):
        return "高", 3
    if any(d in url for d in med):
        return "中", 2
    return "低", 1

# ---------------------------------------------------------------------------
# GOOGLE CUSTOM SEARCH (async wrapper)
# ---------------------------------------------------------------------------

@retry(
    stop=stop_after_attempt(3),
    wait=wait_exponential_jitter(initial=1, max=10),
    retry=retry_if_exception_type(Exception),
)
async def google_search(claim: str) -> List[Dict]:
    """Thread‑off blocking googleapiclient call to keep async loop free."""

    def _sync_search() -> List[Dict]:
        service = build("customsearch", "v1", developerKey=GOOGLE_API_KEY)
        res = (
            service.cse()
            .list(q=claim, cx=CUSTOM_SEARCH_ENGINE_ID, num=10, hl="ja")
            .execute()
        )
        items = res.get("items", [])
        results = []
        for it in items:
            url = it["link"]
            if not url.startswith("https://"):
                continue
            reliability, score = get_source_reliability(url)
            results.append(
                {
                    "title": it["title"],
                    "snippet": it.get("snippet", ""),
                    "link": url,
                    "reliability": reliability,
                    "score": score,
                }
            )
        # スコア順に上位 MAX_SOURCES 件
        return sorted(results, key=lambda x: x["score"], reverse=True)[:MAX_SOURCES]

    return await asyncio.to_thread(_sync_search)

# ---------------------------------------------------------------------------
# OPENAI CALL HELPERS
# ---------------------------------------------------------------------------

@retry(
    stop=stop_after_attempt(3),
    wait=wait_exponential_jitter(initial=1, max=10),
    retry=retry_if_exception_type(OpenAIError),
)
async def openai_chat(messages: List[Dict], *, model: str = "gpt-4o-mini") -> str:
    rsp = await client.chat.completions.create(
        model=model,
        messages=messages,
        temperature=0.0,
    )
    return rsp.choices[0].message.content.strip()


import json
BULLET_SYS = {
    "role": "system",
    "content": (
        "You are a technical writer. Convert each bullet point into a complete, "
        "self-contained factual sentence, using the slide title or parent bullet "
        "as context when needed. Return JSON → {'sentences':[...]}."
    ),
}
BULLET_FUNC_SPEC = {
    "name": "return_sentences",
    "parameters": {
        "type": "object",
        "properties": {
            "sentences": {"type": "array", "items": {"type": "string"}},
        },
        "required": ["sentences"],
    },
}

async def bullets_to_sentences(bullets: list[str]) -> list[str]:
    """`bullets` = 箇条書き1行ごとのリスト"""
    prompt = "\n".join(bullets)
    rsp = await client.chat.completions.create(
        model="gpt-4o-mini",
        messages=[BULLET_SYS, {"role": "user", "content": prompt}],
        functions=[BULLET_FUNC_SPEC],
        function_call={"name": "return_sentences"},
        temperature=0.0,
    )
    args = rsp.choices[0].message.function_call.arguments
    return json.loads(args)["sentences"]


# ---------------------------------------------------------------------------
# CLAIM EXTRACTION (function‑callingで JSON 安定出力)
# ---------------------------------------------------------------------------

EXTRACT_SYS_MSG = {
    "role": "system",
    "content": (
        "あなたは学術エディタです。"
        "入力テキストの中から“事実を述べる文”にあたる箇所を、"
        "一字一句**改変せずに**すべて抜き出し、"
        "JSON で返してください。"
    ),
}
EXTRACT_FUNC_SPEC = {
    "name": "return_claims",
    "description": "テキスト中の主張とその開始・終了位置を返す",
    "parameters": {
        "type": "object",
        "properties": {
            "claims": {
                "type": "array",
                "items": {
                    "type": "object",
                    "properties": {
                        "text": {"type": "string"},
                        "start": {"type": "integer"},
                        "end": {"type": "integer"},
                    },
                    "required": ["text", "start", "end"],
                },
            }
        },
        "required": ["claims"],
    },
}


async def extract_claims(text: str) -> List[str]:
    messages = [EXTRACT_SYS_MSG, {"role": "user", "content": text}]
    rsp = await client.chat.completions.create(
        model="gpt-4o-mini",
        messages=messages,
        functions=[EXTRACT_FUNC_SPEC],
        function_call={"name": "return_claims"},
        temperature=0.0,
    )
    payload = rsp.choices[0].message.function_call.arguments
    import json

    claims = json.loads(payload)["claims"]
    return claims

# ---------------------------------------------------------------------------
# EVIDENCE‑AWARE VERDICT
# ---------------------------------------------------------------------------

async def evidence_based_verdict(claim: str, sources: List[Dict]) -> str:
    # ❶ 証拠を整形（タイトル＋スニペット＋URL）
    evidence_lines = []
    for i, s in enumerate(sources, 1):
        evidence_lines.append(
            f"{i}. [{s['reliability']}] {s['title']} — {s['snippet']} ({s['link']})"
        )
    evidence_txt = "\n".join(evidence_lines) or "〈証拠なし〉"

    # ❷ プロンプトに証拠を埋め込む
    sys_msg = {
        "role": "system",
        "content": (
            "あなたは証拠付きファクトチェッカーです。"
            "必ず証拠のみを根拠に SUPPORTED / REFUTED / NOT SURE のいずれかで判定し、"
            "理由は30字以内で述べてください。"
        ),
    }
    user_msg = {
        "role": "user",
        "content": f"## 主張\n{claim}\n\n## 証拠\n{evidence_txt}",
    }
    return await openai_chat([sys_msg, user_msg], model="gpt-4o")


# ---------------------------------------------------------------------------
# CORE DRIVER (async)
# ---------------------------------------------------------------------------

async def verify_claims(claims: List[str]) -> List[Dict]:
    tasks = []
    for c in claims:
        tasks.append(_process_single_claim(c))
    return await asyncio.gather(*tasks)


async def _process_single_claim(claim: str) -> Dict:
    sources = await google_search(claim)
    verdict = await evidence_based_verdict(claim, sources)
    return {
        "claim": claim,
        "verdict": verdict,
        "sources": sources,
    }

# ---------------------------------------------------------------------------
# DEMO (CLI)
# ---------------------------------------------------------------------------

if __name__ == "__main__":
    parser = argparse.ArgumentParser()
    parser.add_argument("--file", required=True, help="pptx or pdf")
    args = parser.parse_args()

    async def main():
        slide_text, error = read_document(args.file)
        if error:
            print(error)
            sys.exit(1)
        bullets = sanitize_text(slide_text).split("\n")

        # ❷ Bullet → Sentence 変換   ←★ NEW
        sentences = await bullets_to_sentences(bullets)
        claims = await extract_claims("\n".join(sentences))
        print("抽出主張:", claims)


    asyncio.run(main())
