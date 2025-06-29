import os
import re
import pptx
import PyPDF2

def read_document(file_path):
    print(f"'{file_path}' を読み込んでいます...")
    _, extension = os.path.splitext(file_path)
    text = ""
    try:
        if extension == ".pptx":
            prs = pptx.Presentation(file_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
        elif extension == ".pdf":
            with open(file_path, "rb") as f:
                reader = PyPDF2.PdfReader(f)
                for page in reader.pages:
                    text += page.extract_text() + "\n"
        else:
            return None, f"対応していないファイル形式です: {extension}"
        print("ファイルの読み込みが完了しました。")
        return text, None
    except Exception as e:
        return None, f"ファイルの読み込み中にエラーが発生しました: {e}"
    
def extract_lines(slide):
    """(text, level, y_pos, x_pos) を返す"""
    lines = []
    for shape in slide.shapes:
        if not hasattr(shape, "text_frame"):  # 図形・画像などは無視
            continue
        for p in shape.text_frame.paragraphs:
            txt = p.text.strip()
            if not txt:
                continue
            # 箇条書きレベル: 0=タイトル/本文、1=•、2=–– など
            lvl = p.level
            y, x, _, _ = shape.left, shape.top, shape.width, shape.height
            lines.append((txt, lvl, y, x))
    # 画面上: 上→下、左→右
    return sorted(lines, key=lambda t: (t[2], t[3]))


def sanitize_text(text: str) -> str:
    """余分な空白/改行を削除"""
    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r"\n{2,}", "\n", text)
    return text.strip()